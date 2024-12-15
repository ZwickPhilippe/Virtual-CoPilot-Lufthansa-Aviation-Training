import xml.etree.ElementTree as ET
import os
from llama_index.core.schema import TextNode, NodeRelationship, RelatedNodeInfo
from llama_index.core import VectorStoreIndex, PromptTemplate
import re
import json
from llama_index.core import StorageContext, load_index_from_storage
from llama_index.core import QueryBundle
from llama_index.core.schema import NodeWithScore
from llama_index.core.retrievers import BaseRetriever, VectorIndexRetriever
from typing import List
from llama_index.core.postprocessor import SentenceTransformerRerank
from pptx import Presentation
import math


# XML preprocessing & index creation

def append_unit_to_measure_in_text(XML_str: str) -> str:
    # <measure category="Length" show-unit="false" system="Other" unit="NM">2</measure>
    # ->
    # <measure category="Length" show-unit="false" system="Other" unit="NM">2 NM</measure>

    root = ET.fromstring(XML_str)

    for measure_element in root.findall('.//measure'):
        unit = measure_element.get('unit')
        if unit and measure_element.text and not unit in measure_element.text.strip():
            if unit == "FL":  # flight level comes in front
                measure_element.text = "FL " + measure_element.text.strip()
            else:  # all other units go at end
                measure_element.text = measure_element.text.strip() + " " + unit

    return ET.tostring(root, encoding='utf-8').decode('utf-8')


def remove_all_attrs_from_elem(XML_str, input_elements):
    root = ET.fromstring(XML_str)

    for element_name in input_elements:
        for element in root.findall(f'.//{element_name}'):
            element.attrib.clear()

    return ET.tostring(root, encoding='utf-8').decode('utf-8')


def remove_elem_leave_text(XML_str, clean_list):
    XML_str = remove_all_attrs_from_elem(XML_str, clean_list)
    for i in clean_list:
        XML_str = XML_str.replace(f'<{i}>', ' ')
        # if at end of sentence we don't want space
        XML_str = XML_str.replace(f'</{i}>.', '.')
        XML_str = XML_str.replace(f'</{i}>', ' ')
    return XML_str


def remove_fragments_from_text(XML_str, fragments):
    for i in fragments:
        XML_str = XML_str.replace(f'{i}', ' ')
    return XML_str.replace('  ', ' ')


def remove_attrs_from_all_elems(XML_str, attrs):
    root = ET.fromstring(XML_str)

    for elem in root.iter():
        for attr in attrs:
            if attr in elem.attrib:
                del elem.attrib[attr]
    return ET.tostring(root, encoding='utf-8').decode('utf-8')


def parse_titles_bottom_up(XML_structure_path: str) -> dict:
    """
    Use for FCOM.

    :param XML_structure_path: Path to the XML file containing structure of individual XML files.
    :return: Dictionary mapping XML filenames to a tuple of the title and their position number.
    """
    tree = ET.parse(XML_structure_path)
    root = tree.getroot()
    title_dict = {}
    counter = 0

    # Helper function to traverse XML tree
    def traverse_path(element, current_titles):
        nonlocal counter
        for child in element:
            # Add title if it's part of the path to <du-sol>
            if child.tag in {"psl", "du-inv", "du-sol"}:
                title_element = child.find("title")
                if title_element is not None and title_element.text is not None:
                    # Add title to the current path
                    current_titles.append(title_element.text.strip())

            # Check if we've reached a <du-sol> with <sol-content>
            if child.tag == "du-sol":
                sol_title = ", ".join(current_titles)
                href = child.find("sol-content-ref").get("href")
                filename = href.split('/')[-1]
                title_dict[filename] = (sol_title, counter)
                counter += 1
            else:
                # Recursively check further children only if we haven't reached <du-sol>
                # Pass a copy to prevent modification
                traverse_path(child, current_titles[:])

            # Remove the last title after processing this level to reset the path
            if child.tag in {"psl", "du-inv", "du-sol"}:
                current_titles.pop()  # Clean up after this path

    # Start the traversal
    traverse_path(root, [])
    return title_dict


def parse_titles_bottom_up_PEGMA(XML_structure_path: str) -> dict:
    """
    Use for FCTM and FCM.

    :param XML_structure_path: Path to the XML file containing structure of individual XML files.
    :return: Dictionary mapping XML filenames to a tuple of the title and their position number.
    """
    tree = ET.parse(XML_structure_path)
    root = tree.getroot()
    title_dict = {}
    counter = 0

    # Helper function to traverse XML tree
    def traverse_path(element, current_titles):
        nonlocal counter
        for child in element:
            # Add title if it's part of the path to <solution>
            if child.tag in {"information-node", "invariant", "solution"}:
                title = child.get("title")
                if title is not None:
                    # Add title to the current path
                    current_titles.append(title.strip())

            # Check if we've reached a <solution> with <content>
            if child.tag == "solution":
                sol_title = ", ".join(current_titles)
                href = child.find("content").get("href")
                filename = href.split('/')[-1]
                title_dict[filename] = (sol_title, counter)
                counter += 1
            else:
                # Recursively check further children only if we haven't reached <solution>
                # Pass a copy to prevent modification
                traverse_path(child, current_titles[:])

            # Remove the last title after processing this level to reset the path
            if child.tag in {"information-node", "invariant", "solution"}:
                current_titles.pop()  # Clean up after this path

    # Start the traversal
    traverse_path(root, [])
    return title_dict


def clean_XML_txt(txt: str, XML_full_clean: bool) -> str:
    """
    Top-level cleaning function

    :param txt: XML file content for one file
    :param XML_full_clean: indicator for full clean
    :return: Cleaned XML file content
    """
    txt = append_unit_to_measure_in_text(txt)
    # txt = remove_all_attrs_from_elem(txt, ['para']) #superseded by remove_attrs_from_all_elems()
    txt = remove_attrs_from_all_elems(txt,
                                      ['layer', 'key', 'code', 'lid', 'isl-id', 'checksum', 'id', 'lid', 'xmlns:xsi',
                                       'xsi:noNamespaceSchemaLocation'])
    txt = remove_elem_leave_text(txt, ['abb', 'measure', 'tech-label'])
    txt = remove_fragments_from_text(txt, ['<?xml version="1.0" encoding="UTF-8"?>',
                                           '<!DOCTYPE description PUBLIC ',
                                           '"-//AIRBUS OPS//DTD DESCRIPTION_EXCH-VER10.0//EN" "../DTD/description-exch.dtd">',
                                           'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance"',
                                           'xsi:noNamespaceSchemaLocation="../DTD/DU_l1item.xsd"',
                                           ])
    txt = txt.replace('\n', ' ')
    txt = txt.replace('  ', ' ')

    if XML_full_clean:
        txt = re.sub(r'<[^>]*>', ' ',
                     txt)  # remove_all_elements: Matches any opening tag "<" followed by non-">" characters "*" times
        # replace "  " (only those!) after above line
        txt = txt.replace('  ', ' ')

    return txt.strip()


def extract_du_sol_titles(xml_file_path):
    # Parse the XML file
    tree = ET.parse(xml_file_path)
    root = tree.getroot()

    # Initialize a list to store titles
    titles = []

    # Get the very first <title> in the XML (outside of <du-sol>)
    first_title = root.find(".//title")
    if first_title is not None:
        titles.append(first_title.text)

    # Find all <du-sol> elements and extract their <title>
    du_sol_elements = root.findall(".//du-sol")
    for du_sol in du_sol_elements:
        du_sol_title = du_sol.find("title")
        if du_sol_title is not None and du_sol_title.text not in titles:
            titles.append(du_sol_title.text)

    return titles


def extract_information_node_titles(xml_file_path):
    # Parse the XML file
    tree = ET.parse(xml_file_path)
    root = tree.getroot()

    # Initialize a list to store titles
    titles = []

    # Find all <invariant> elements and extract their 'title' attributes
    invariant_elements = root.findall(".//information-node")
    for invariant in invariant_elements:
        title = invariant.get("title")  # Get the 'title' attribute
        if title and title not in titles:  # Check if 'title' attribute exists
            titles.append(title)

    return titles


def create_nodes_with_relationships(XML_folder_path: str, file_title_relation: dict, XML_full_clean: bool) -> list:
    """
    Use for fancy index creation: create nodes & enhance with relationships.

    :param XML_folder_path: Path to folder with all individual XML files.
    :param file_title_relation: dict from parse_titles_bottom_up(_PEGMA)()
    :param XML_full_clean: indicator for full clean in clean_XML_txt()
    :return: Dictionary mapping XML filenames to a tuple of the title and their position number. 
    """
    # CREATE NODES
    nodes = []
    for filename in os.listdir(XML_folder_path):
        if filename.endswith(".xml"):
            file_path = os.path.join(XML_folder_path, filename)
            with open(file_path, "r", encoding="utf-8") as f:
                txt = f.read()

            try:
                file_title = file_title_relation.get(filename)[0]
            except:
                # file not in master structure => don't include it in our index either!
                print(f"File not part of overall structure: {filename}")
                continue

            txt = clean_XML_txt(txt, XML_full_clean)

            txt = f"Title: {file_title}\n{txt}"
            node = TextNode(text=txt, id_=filename)  # filename is unique
            nodes.append(node)

    # RELATIONSHIPS
    nodes.sort(key=lambda node: file_title_relation.get(
        node.id_)[1])  # sort nodes by index

    for i in range(len(nodes) - 1):
        current_node = nodes[i]
        next_node = nodes[i + 1]

        # Ensure consistent relationship direction based on filename order
        current_node.relationships[NodeRelationship.NEXT] = RelatedNodeInfo(
            node_id=next_node.id_
        )
        next_node.relationships[NodeRelationship.PREVIOUS] = RelatedNodeInfo(
            node_id=current_node.id_
        )

    return nodes


def create_index(XML_folder_path: str, file_title_relation: dict = None,
                 XML_full_clean: bool = False) -> VectorStoreIndex:
    """
    Top-level function, to be called outside of my_XML.py

    :param XML_folder_path: Path to folder with all individual XML files.
    :param file_title_relation: dict from parse_titles_bottom_up(_PEGMA)()
    :param XML_full_clean: indicator for full clean in clean_XML_txt()
    :return: Vector Store for this XML document
    """
    if (file_title_relation):
        nodes = create_nodes_with_relationships(
            XML_folder_path, file_title_relation, XML_full_clean)
    else:  # nodes have no order, no relationships, no titles
        nodes = []
        for file_name in os.listdir(XML_folder_path):
            if file_name.endswith('.xml'):
                file_path = os.path.join(XML_folder_path, file_name)
                with open(file_path, 'r', encoding='utf-8') as f:
                    txt = f.read()
                txt = clean_XML_txt(txt, XML_full_clean)
                node = TextNode(text=txt, id_=file_name)
                nodes.append(node)

    return VectorStoreIndex(nodes=nodes, show_progress=True)


def merge_all_duSol(XML_structure_path: str, XML_folder_path: str, output_folder) -> None:
    """
    EXPERIMENTAL
    For FCOM, merge all du-sol (children docs) of the same du-inf (parent docs)

    :param XML_structure_path: Path to the XML file containing structure of individual XML files.
    :param XML_folder_path: Path to folder with all individual XML files.
    :param output_folder: save location of newly merged XML files
    :return: Vector Store for this XML document
    """
    root = ET.parse(XML_structure_path).getroot()

    for elem in root.findall(".//du-inv"):  # iterate over 4829 elems
        output_file = os.path.join(output_folder, elem.get('code') + ".xml")

        newfile = ET.Element("myRoot")

        for sol in elem.findall(".//du-sol"):
            sol_path = os.path.join(XML_folder_path, sol.get('code') + ".xml")
            sol_root = ET.parse(sol_path).getroot()
            newfile.append(sol_root)

        new_tree = ET.ElementTree(newfile)
        new_tree.write(output_file)


def parse_titles_bottom_up_duInv(XML_structure_path: str) -> dict:
    """
    EXPERIMENTAL
    Use on merge_all_duSol() XML files.

    :param XML_structure_path: Path to the XML file containing structure of individual XML files.
    :return: Dictionary mapping XML filenames to a tuple of the title and their position number.
    """
    tree = ET.parse(XML_structure_path)
    root = tree.getroot()
    title_dict = {}
    counter = 0

    # Helper function to traverse XML tree
    def traverse_path(element, current_titles):
        nonlocal counter
        for child in element:
            # Add title if it's part of the path to <du-sol>
            if child.tag in {"psl", "du-inv"}:
                title_element = child.find("title")
                if title_element is not None and title_element.text is not None:
                    # Add title to the current path
                    current_titles.append(title_element.text.strip())

            # Check if we've reached a <du-inv>
            if child.tag == "du-inv":
                inv_title = ", ".join(current_titles)
                href = child.get("code") + ".xml"
                title_dict[href] = (inv_title, counter)
                counter += 1
            else:
                # Recursively check further children only if we haven't reached <du-inv>
                # Pass a copy to prevent modification
                traverse_path(child, current_titles[:])

            # Remove the last title after processing this level to reset the path
            if child.tag in {"psl", "du-inv"}:
                current_titles.pop()  # Clean up after this path

    # Start the traversal
    traverse_path(root, [])
    return title_dict


# Functions used in the streamlit
def get_id_lookup_dict_for_final_title_print(base_path_MANUALS):
    lookup_dict = {}
    my_path = base_path_MANUALS + "/storage/dict_title_filename_document"

    for filename in os.listdir(my_path):
        with open(os.path.join(my_path, filename), 'r') as f:
            data = json.load(f)
        if len(set(lookup_dict.keys()).intersection(set(data.keys()))) > 0:
            assert False, "keys should be unqiue"
        else:
            lookup_dict.update(data)
    return lookup_dict


def dict_to_markdown_table(data):
    columns = list(data.keys())
    rows = list(zip(*data.values()))

    # Construct the Markdown table
    header = "| " + " | ".join(columns) + " |"
    separator = "| " + " | ".join(["---"] * len(columns)) + " |"
    rows_md = "\n".join("| " + " | ".join(row) + " |" for row in rows)

    return f"{header}\n{separator}\n{rows_md}"


def load_indeces_retriever(indeces_topk: dict) -> list:
    option = "_llama3B_bgeL"
    res = []
    for i, k in indeces_topk:
        storage_context = StorageContext.from_defaults(
            persist_dir="./_MANUALS/storage/" + i + option
        )
        index = load_index_from_storage(storage_context)
        res.append(VectorIndexRetriever(index=index, similarity_top_k=k))
    return res


def my_retriever_builder(indices_topk, llm, reranker=None):
    indices = load_indeces_retriever(indices_topk)
    custom_retriever = CustomRetriever(indices, llm, reranker)
    return custom_retriever


class CustomRetriever(BaseRetriever):
    """Looks up the query in all indices and concatenates the retrieved nodes to one list."""

    def __init__(
            self,
            vector_retrievers: List[VectorIndexRetriever],
            llm,
            my_reranker: SentenceTransformerRerank
    ) -> None:
        """Init params."""
        self._vector_retrievers = vector_retrievers
        self.my_reranker = my_reranker
        self.llm = llm
        self.applicable_index_retrieval_template = (
            "The prompt you will receive is a question,that should be aircraft related. In the prompt, there are most "
            "likely gonna be aircraft models included. Typical aircraft models include:"
            "---------------------\n"
            "{context_str}"
            "\n---------------------\n"
            "Please identify if the question is to any extend aircraft related. If this is the case, return 'True', "
            "if not return 'False'. It is important that you respond with either True or False and really only return "
            "True, if is actually related to the query."
            "The query is: {query_str}. Explain your answer please.\n"
        )
        self.contex_aircraft_models = 'A320, A321, A322, A380, Airbus, Boeing'
        super().__init__()

    def my_rerank(self, vector_nodes, query_bundle):
        return self.my_reranker.postprocess_nodes(nodes=vector_nodes, query_bundle=query_bundle)

    def find_applicable_index(self, query_bundle):
        prompt_template = PromptTemplate(
            self.applicable_index_retrieval_template)
        message = prompt_template.format_messages(context_str=self.contex_aircraft_models,
                                                  query_str=query_bundle.query_str)

        response = self.llm.chat(message)

        if response is None or 'True' in response.message.content:
            print("Related to the question")
            return self._vector_retrievers
        else:
            print("Not related to the question")
            return []

    def _retrieve(self, query_bundle: QueryBundle) -> List[NodeWithScore]:
        """Retrieve nodes given query."""

        print(f"Retrieval for following query: {query_bundle.query_str}")

        applic_idx = self.find_applicable_index(query_bundle)

        vector_nodes = []
        for idx in applic_idx:
            print("Retrieval on index-------------")
            vector_nodes += idx.retrieve(query_bundle)

        if self.my_reranker:
            print(
                f"BEFORE reranking we have following {len(vector_nodes)} retrieved nodes -----------------------")
            print(vector_nodes)

            vector_nodes = self.my_rerank(vector_nodes, query_bundle)

            print(
                f"AFTER reranking we have following {len(vector_nodes)} retrieved nodes -----------------------")
            print(vector_nodes)

        # assert False, "remove this if you want the LLM to get queried with the vector nodes to draft an answer"

        return vector_nodes


def extract_text_from_slide(slide):
    title = None
    content = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text.strip()

        # If the shape has a title-like content, assign it as the slide's title
        if shape == slide.shapes.title:
            title = text
        else:
            content.append(text)

    return title, content


def extract_content_from_pptx(file_path):
    prs = Presentation(file_path)
    slides_content = []

    for i, slide in enumerate(prs.slides):
        title, content = extract_text_from_slide(slide)
        slide_data = {
            "slide_number": i + 1,
            "title": title if title else "Untitled Slide",
            "content": content
        }
        slides_content.append(slide_data)

    return slides_content


def process_pptx_files(folder_path, json_output_path):
    all_presentations = {}

    for file_name in os.listdir(folder_path):
        if file_name.lower().endswith('.pptx'):
            file_path = os.path.join(folder_path, file_name)
            print(f"Processing file: {file_name}")
            slides_content = extract_content_from_pptx(file_path)

            # Use the file name (without extension) as the topic
            topic = os.path.splitext(file_name)[0]
            all_presentations[topic] = slides_content

    with open(json_output_path, 'w', encoding='utf-8') as f:
        json.dump(all_presentations, f, indent=4, ensure_ascii=False)


def briefing_JSON_split_evenly(data, max_chars=4000):
    # Could also be adapted for FCOM and other docs in the future.

    total_chars = len(str(data))

    # If the data is within the limit, no splitting is needed
    if total_chars < max_chars:
        return [data]

    # Calculate the number of chunks
    num_chunks = math.ceil(total_chars / max_chars)
    max_chunk_size = math.ceil(total_chars / num_chunks)
    print(num_chunks)
    x = "slkdjf"

    chunks = [[]]
    cur_chunk_size = 0
    index = 0

    for page in data:
        if cur_chunk_size < max_chunk_size:
            chunks[index].append(page)
        else:
            cur_chunk_size = 0
            index = index + 1
            chunks.append([])
            chunks[index].append(page)
        cur_chunk_size += len(str(page))

    return chunks


def create_chunks_for_briefings_json(extracted_briefings, output_directory):
    # Splits chunks evenly

    with open(extracted_briefings, 'r', encoding='utf-8') as f:
        x = json.load(f)

    for key, value in x.items():
        json_chunks = briefing_JSON_split_evenly(value)

        # Save each chunk to a separate file
        for i, chunk in enumerate(json_chunks):
            suffix = f"_part{i + 1}" if len(json_chunks) > 1 else ""
            data_to_save = {key: chunk}
            os.makedirs(output_directory, exist_ok=True)
            output_file = os.path.join(output_directory, f"{key}{suffix}.json")
            with open(output_file, 'w', encoding='utf-8') as out_file:
                json.dump(data_to_save, out_file, indent=4, ensure_ascii=False)


def create_simple_chunks_for_briefings_json(extracted_briefings, output_directory):
    with open(extracted_briefings, 'r', encoding='utf-8') as f:
        x = json.load(f)

    for key, value in x.items():
        data_to_save = {key: value}
        os.makedirs(output_directory, exist_ok=True)
        output_file = os.path.join(output_directory, f"{key}.json")
        with open(output_file, 'w', encoding='utf-8') as out_file:
            json.dump(data_to_save, out_file, indent=4, ensure_ascii=False)
